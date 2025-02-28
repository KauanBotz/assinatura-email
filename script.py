import os
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import win32com.client
import re

pasta_base = r"\\Arquivos - Assinatura"
pasta_assinatura = os.path.join(pasta_base, "assinatura")

print("Bem-vindo ao Gerador de Assinaturas.")
print("Construtora Sant'Anna")
print("A assinatura será salva na pasta (Arquivos-Assinatura), localizada no diretório (assinatura).")
print("")

if not os.path.exists(pasta_assinatura):
    os.makedirs(pasta_assinatura)

def editar_assinatura(pptx_path, nome, funcao, numero=None, endereco=None):
    if not os.path.exists(pptx_path):
        print(f"Erro: O arquivo '{pptx_path}' não foi encontrado!")
        return
    
    print(f"Arquivo '{pptx_path}' carregado com sucesso!")

    prs = Presentation(pptx_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if '{{Nome}}' in shape.text:
                    shape.text = shape.text.replace('{{Nome}}', nome)
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0x59, 0x59, 0x59)  # Cor #595959
                            run.font.size = Pt(28)
                if '{{Funcao}}' in shape.text:
                    shape.text = shape.text.replace('{{Funcao}}', funcao)
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0x59, 0x59, 0x59)  # Cor #595959
                if '{{Numero}}' in shape.text and numero:
                    shape.text = shape.text.replace('{{Numero}}', numero)
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0x00, 0xA5, 0x51)  # Cor #3A8C62
                elif '{{Numero}}' in shape.text:
                    shape.text = shape.text.replace('{{Numero}}', '')
                if '{{Endereço}}' in shape.text and endereco:
                    shape.text = shape.text.replace('{{Endereço}}', endereco)
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0x00, 0xA5, 0x51)  # Cor #3A8C62
                elif '{{Endereço}}' in shape.text:
                    shape.text = shape.text.replace('{{Endereço}}', '')

    temp_pptx_path = os.path.join(pasta_assinatura, 'temp_assinatura.pptx')
    prs.save(temp_pptx_path)  
    print(f"PPTX modificado salvo como '{temp_pptx_path}'.")
    
    jpg_path = os.path.join(pasta_assinatura, f'{nome}.jpg')
    converter_para_jpg(temp_pptx_path, jpg_path)

def converter_para_jpg(pptx_path, jpg_path):
    ppt = None
    try:
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = True 
        presentation = ppt.Presentations.Open(os.path.abspath(pptx_path))

        presentation.SaveAs(os.path.abspath(jpg_path), 18)
        presentation.Close()
        print(f"Arquivo JPG gerado e salvo como '{jpg_path}'.")

    except Exception as e:
        print(f"Erro ao converter para JPG: {e}")
    finally:
        if ppt:
            ppt.Quit()

def escolher_endereco():
    print("Escolha o endereço:")
    print("(1) Rua São Pedro da Aldeia, 1200 - Pilar")
    print("(2) Governador Valadares")
    print("(3) Conselheiro Pena")
    print("(4) Digite um endereço")
    print("")
    escolha = input("Digite o número da opção desejada: ")
    print("Endereço selecionado")
    print("")
    
    enderecos = {
        "1": "Rua São Pedro da Aldeia, 1200 - Pilar",
        "2": "Governador Valadares",
        "3": "Conselheiro Pena"
    }
    
    if escolha == "4":
        endereco_personalizado = input("Digite o endereço: ")
        return endereco_personalizado
    else:
        return enderecos.get(escolha, "Endereço não especificado")

import re

def validar_numero(numero):
    padrao = r"^\(\d{2}\) 9?\d{4}-\d{4}$"
    if re.match(padrao, numero):
        return True
    else:
        print("Erro: O número de telefone não está no formato correto. Exemplo: (31) 9XXXX-XXXX ou (31) XXXX-XXXX")
        return False

while True:
    nome = input("Digite o nome: ")
    print("")
    funcao = input("Digite a função: ")
    print("")
    
    while True:
        numero = input("Digite o seu número de telefone (Exemplo: (31) XXXX-XXXX ou (31) 9XXXX-XXXX): ")
        if numero and validar_numero(numero):
            break
    print("")

    endereco = escolher_endereco()

    os.chdir(r"C:\Users\aprendiz.ti\Desktop\Arquivos - Assinatura")

    editar_assinatura(
        pptx_path = os.path.join(pasta_base, "modelo-assinaturaNovo.pptx"),
        nome=nome,
        funcao=funcao,
        numero=numero if numero else None,
        endereco=endereco
    )
    
    print("")
    continuar = input("Deseja gerar outra assinatura? (sim/não): ")
    if continuar.lower() != 'sim':
        break
    print("")
    print("Gerando um nova assinatura")

